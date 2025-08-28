
import os
import re
from typing import List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from .schemas import SlideContent

# --- THEMES for better contrast and thematic feel ---
THEMES: Dict[str, Dict[str, RGBColor]] = {
    "Minimalist_Dark": {
        "background": RGBColor(18, 18, 18),      # Near-black
        "primary": RGBColor(240, 240, 240),      # Off-white for text
        "secondary": RGBColor(0, 150, 255),      # A vibrant but not jarring blue
        "accent": RGBColor(130, 130, 130)        # Lighter gray for subtitles
    },
    "Business_Corporate": {
        "background": RGBColor(245, 245, 245),   # Off-white/light gray
        "primary": RGBColor(10, 30, 60),         # Professional navy blue text
        "secondary": RGBColor(0, 123, 255),      # Standard corporate blue
        "accent": RGBColor(108, 117, 125)        # Medium gray
    },
    "Education_Creative": {
        "background": RGBColor(255, 253, 248),   # Very light cream
        "primary": RGBColor(50, 50, 50),         # Dark charcoal text
        "secondary": RGBColor(255, 110, 60),     # Engaging orange accent
        "accent": RGBColor(120, 120, 120)        # Soft gray for subtitles
    },
    "Historical_Vintage": {
        "background": RGBColor(245, 245, 220),   # Beige/Parchment
        "primary": RGBColor(80, 50, 20),         # Dark, sepia-toned brown
        "secondary": RGBColor(139, 69, 19),      # Saddle brown accent
        "accent": RGBColor(160, 82, 45)          # Lighter brown
    },
    "Technology_Futuristic": {
        "background": RGBColor(10, 0, 25),       # Deep dark purple/blue
        "primary": RGBColor(0, 255, 127),        # Spring green text ("hacker green")
        "secondary": RGBColor(190, 70, 255),     # Electric purple/magenta
        "accent": RGBColor(105, 105, 105)        # Dim gray
    },
    "Environmental_Natural": {
        "background": RGBColor(240, 255, 240),   # Honeydew/pale green
        "primary": RGBColor(0, 80, 0),           # Deep, earthy green
        "secondary": RGBColor(34, 139, 34),      # Forest green
        "accent": RGBColor(107, 142, 35)         # Olive
    },
    "default": {
        "background": RGBColor(255, 255, 255),
        "primary": RGBColor(0, 0, 0),
        "secondary": RGBColor(100, 100, 100),
        "accent": RGBColor(150, 150, 150)
    },
}

# --- Thematic Fonts for deeper theme integration ---
THEME_FONTS: Dict[str, Dict[str, str]] = {
    "Minimalist_Dark": {"title": "Segoe UI", "body": "Calibri"},
    "Business_Corporate": {"title": "Arial", "body": "Helvetica"},
    "Education_Creative": {"title": "Century Gothic", "body": "Trebuchet MS"},
    "Historical_Vintage": {"title": "Garamond", "body": "Georgia"},
    "Technology_Futuristic": {"title": "Tw Cen MT", "body": "Verdana"},
    "Environmental_Natural": {"title": "Rockwell", "body": "Calibri"},
    "default": {"title": "Segoe UI", "body": "Calibri"},
}

OUTPUT_DIR = "outputs"
os.makedirs(OUTPUT_DIR, exist_ok=True)


def add_background_color(slide, color: RGBColor):
    """Add a solid background color to a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def create_presentation_file(topic: str, presenter_name: str, template_name: str, slide_contents: List[SlideContent]) -> str:
    # --- Get Theme Colors and Fonts ---
    theme = THEMES.get(template_name, THEMES["default"])
    fonts = THEME_FONTS.get(template_name, THEME_FONTS["default"])
    
    background_color = theme["background"]
    primary_color = theme["primary"]
    secondary_color = theme["secondary"]
    accent_color = theme["accent"]
    title_font = fonts["title"]
    body_font = fonts["body"]

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)  # 16:9 Aspect Ratio

    # --- Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, background_color)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2.5))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf_title.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.2
    
    title_text = slide_contents[0].title if slide_contents else topic
    run = p.add_run()
    run.text = title_text[:60] + "..." if len(title_text) > 60 else title_text
    run.font.name = title_font
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = primary_color

    subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(11.63), Inches(1))
    tf_sub = subtitle_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.text = f"Presented by {presenter_name}"
    p_sub.font.name = body_font
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = accent_color

    # --- Content Slides ---
    for i, slide_data in enumerate(slide_contents[1:], start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background_color(slide, background_color)

        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
        tf_header = header_box.text_frame
        p_head = tf_header.paragraphs[0]
        p_head.text = topic.upper()
        p_head.font.name = body_font
        p_head.font.size = Pt(12)
        p_head.font.bold = True
        p_head.font.color.rgb = accent_color

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p1 = tf_title.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        
        title_parts = re.split(r'[:\-–]', slide_data.title, 1)
        run1 = p1.add_run()
        run1.text = title_parts[0].strip()
        run1.font.name = title_font
        run1.font.size = Pt(28)
        run1.font.bold = True
        run1.font.color.rgb = primary_color
        
        if len(title_parts) > 1:
            run2 = p1.add_run()
            run2.text = f" - {title_parts[1].strip()}"
            run2.font.name = title_font
            run2.font.size = Pt(28)
            run2.font.bold = True
            run2.font.color.rgb = secondary_color

        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(4.4))
        tf_body = body_box.text_frame
        tf_body.word_wrap = True

        for j, bullet_text in enumerate(slide_data.bullets):
            p = tf_body.add_paragraph() if j > 0 else tf_body.paragraphs[0]
            p.line_spacing = 1.3
            p.space_after = Pt(12)
            p.level = 1 if bullet_text.startswith("  ") else 0
            
            clean_text = bullet_text.strip()
            segments = re.split(r'(\*\*.*?\*\*|__.*?__)', clean_text)
            
            for seg in segments:
                run = p.add_run()
                run.font.name = body_font
                run.font.size = Pt(22)
                run.font.color.rgb = primary_color
                
                if seg.startswith("**") and seg.endswith("**"):
                    run.text = seg[2:-2]
                    run.font.bold = True
                    run.font.color.rgb = secondary_color
                elif seg.startswith("__") and seg.endswith("__"):
                    run.text = seg[2:-2]
                    run.font.underline = True
                else:
                    run.text = seg

        if hasattr(slide_data, "references") and slide_data.references:
            ref_box_top = Inches(6.2) 
            ref_box = slide.shapes.add_textbox(Inches(0.8), ref_box_top, Inches(8.2), Inches(0.8))
            p_ref = ref_box.text_frame.paragraphs[0]
            p_ref.text = "Sources: "
            p_ref.font.name = body_font
            p_ref.font.size = Pt(9)
            p_ref.font.bold = True
            p_ref.font.color.rgb = accent_color
            
            for idx, url in enumerate(slide_data.references[:2]):
                if idx > 0:
                    p_ref.add_run().text = " | "
                
                run = p_ref.add_run()
                run.text = url[:40] + "..." if len(url) > 40 else url
                run.hyperlink.address = url
                run.font.size = Pt(9)
                run.font.color.rgb = accent_color

        page_widget_top = Inches(6.6) 
        page_widget = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), page_widget_top, Inches(1.8), Inches(0.4))
        fill = page_widget.fill
        fill.solid()
        fill.fore_color.rgb = secondary_color
        page_widget.line.fill.background()
        
        p_page = page_widget.text_frame.paragraphs[0]
        p_page.alignment = PP_ALIGN.CENTER
        run_page = p_page.add_run()
        run_page.text = f"Page No. {i+1:02d}"
        run_page.font.name = body_font
        run_page.font.size = Pt(11)
        run_page.font.bold = True
        
        is_dark_bg = sum(secondary_color) < 382 
        run_page.font.color.rgb = RGBColor(255, 255, 255) if is_dark_bg else RGBColor(0,0,0)


    # --- Thank You Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_color(slide, background_color)
    
    thank_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
    p = thank_box.text_frame.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = title_font
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = secondary_color

    questions_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11.33), Inches(1))
    p_q = questions_box.text_frame.paragraphs[0]
    p_q.text = "Questions & Discussion"
    p_q.alignment = PP_ALIGN.CENTER
    p_q.font.name = body_font
    p_q.font.size = Pt(20)
    p_q.font.color.rgb = accent_color

    safe_topic = "".join(x for x in topic if x.isalnum() or x in " -_").rstrip()
    file_path = os.path.join(OUTPUT_DIR, f"{safe_topic}_{template_name}.pptx")
    prs.save(file_path)
    return file_path